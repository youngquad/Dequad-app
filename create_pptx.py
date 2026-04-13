from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Create presentation (widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
BRAND_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0x63, 0x66, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

SCREENSHOT_DIR = "/tmp/ppt_screenshots"
LOGO_PATH = "/tmp/dequad_logo.png"


def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, DARK_BG)
    bg.line.fill.background()


def add_text(slide, left, top, width, height, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


# ===== SLIDE 1: Title Slide =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.4), Inches(1.0), width=Inches(2.5))

add_text(slide, Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.2),
         "DEQUAD", size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(4.7), Inches(12.333), Inches(0.8),
         "AI-Powered Student Wellbeing Platform", size=28, color=BRAND_BLUE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5),
         "Connecting students. Protecting wellbeing. Empowering universities.", size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 2: Core Features =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text(slide, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8),
         "Core Features", size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.667), Inches(1.2), Inches(2), Inches(0.06))
set_shape_fill(line, BRAND_BLUE)
line.line.fill.background()

left_features = [
    "Google OAuth Authentication",
    "Daily Mood Tracking (1-10 scale)",
    "AI-Powered Risk Analysis",
    "Hinge-Style Student Matching",
    "Comment on Likes Feature",
    "End-to-End Encrypted Chat",
]
right_features = [
    "Safeguarding Alerts System",
    "Super Admin Dashboard",
    "University Admin Dashboard",
    "Stripe Premium Subscriptions",
    "Data Export (CSV)",
    "Email Notifications",
]

left_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.5), Inches(5))
tf = left_box.text_frame
tf.word_wrap = True
for i, f in enumerate(left_features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {f}"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.space_after = Pt(16)

right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.5), Inches(5))
tf = right_box.text_frame
tf.word_wrap = True
for i, f in enumerate(right_features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {f}"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.space_after = Pt(16)

# ===== SLIDE 3: Landing Page =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7),
         "Landing Page", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.4),
         "Clean, modern landing with Google Sign-In and university subscription entry point", size=16, color=GRAY)
img_path = os.path.join(SCREENSHOT_DIR, "ppt_01_landing.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.333))

# ===== SLIDE 4: University Subscription =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7),
         "University Subscription", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.4),
         "Universities subscribe at GBP 49.99/month to monitor student wellbeing with full dashboard access", size=16, color=GRAY)
img_path = os.path.join(SCREENSHOT_DIR, "ppt_02_uni_subscribe.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.333))

# ===== SLIDE 5: University Admin Login =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7),
         "University Admin Login", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.4),
         "Dedicated secure login portal for university administrators with auto-generated credentials", size=16, color=GRAY)
img_path = os.path.join(SCREENSHOT_DIR, "ppt_03_uni_login.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.333))

# ===== SLIDE 6: University Admin Dashboard =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7),
         "University Admin Dashboard", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.4),
         "Real-time analytics: student count, mood trends, safeguarding alerts, and data export", size=16, color=GRAY)
img_path = os.path.join(SCREENSHOT_DIR, "ppt_04_dashboard_overview.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.333))

# ===== SLIDE 7: Super Admin Dashboard =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7),
         "Super Admin Dashboard", size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.4),
         "Platform-wide analytics: all students, premium users, safeguarding alerts, and AI insights", size=16, color=GRAY)
img_path = os.path.join(SCREENSHOT_DIR, "ppt_06_super_admin.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(12.333))

# ===== SLIDE 8: Pricing =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
         "Pricing Plans", size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Student plan box
box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(4.5), Inches(5))
set_shape_fill(box1, RGBColor(0x1E, 0x29, 0x3B))
box1.line.color.rgb = BRAND_BLUE
box1.line.width = Pt(2)

add_text(slide, Inches(1.8), Inches(1.8), Inches(4), Inches(0.5),
         "Student Premium", size=26, color=BRAND_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.8), Inches(2.4), Inches(4), Inches(0.7),
         "GBP 4.99/month", size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

student_features = ["Unlimited matches", "Priority support", "Advanced AI insights", "No ads"]
tf_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.3), Inches(3.8), Inches(3))
tf = tf_box.text_frame
tf.word_wrap = True
for i, f in enumerate(student_features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {f}"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.space_after = Pt(10)

# University plan box
box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.333), Inches(1.5), Inches(4.5), Inches(5))
set_shape_fill(box2, RGBColor(0x1E, 0x29, 0x3B))
box2.line.color.rgb = PURPLE
box2.line.width = Pt(2)

add_text(slide, Inches(7.633), Inches(1.8), Inches(4), Inches(0.5),
         "University Dashboard", size=26, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.633), Inches(2.4), Inches(4), Inches(0.7),
         "GBP 49.99/month", size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

uni_features = ["Full dashboard access", "Student monitoring", "Safeguarding alerts", "Data export", "Priority support"]
tf_box = slide.shapes.add_textbox(Inches(7.833), Inches(3.3), Inches(3.8), Inches(3))
tf = tf_box.text_frame
tf.word_wrap = True
for i, f in enumerate(uni_features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {f}"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.space_after = Pt(10)

# ===== SLIDE 9: Demo Users =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8),
         "Demo User Profiles", size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.667), Inches(1.1), Inches(2), Inches(0.06))
set_shape_fill(line, ACCENT_GREEN)
line.line.fill.background()

users = [
    ("Emma Wilson", "Computer Science", "University of Manchester", "Visual Learner"),
    ("James Chen", "Data Science", "University of Birmingham", "Night Owl"),
    ("Sofia Martinez", "Psychology", "University of Leeds", "Early Bird"),
    ("Alex Thompson", "Environmental Science", "University of Bristol", "Group Study"),
    ("Priya Patel", "Business Analytics", "University of Warwick", "Visual Learner"),
    ("Oliver Wright", "Mechanical Engineering", "University of Manchester", "Night Owl"),
    ("Amara Okafor", "Biomedical Science", "University of Manchester", "Early Bird"),
    ("Lucas Fernandez", "Architecture", "University of Leeds", "Visual Learner"),
    ("Zara Ahmed", "Law", "University of Birmingham", "Group Study"),
    ("Ethan Kim", "Mathematics", "University of Warwick", "Solo Study"),
    ("Chloe Williams", "English Literature", "University of Bristol", "Early Bird"),
    ("Daniel Johnson", "Music Production", "University of Manchester", "Night Owl"),
]

# Display in 3 columns of 4
cols = 3
rows = 4
col_width = 4.0
row_height = 1.2
start_x = 0.8
start_y = 1.5

for idx, (name, course, uni, style) in enumerate(users):
    col = idx % cols
    row = idx // cols
    x = start_x + col * col_width
    y = start_y + row * row_height

    # Name
    add_text(slide, Inches(x), Inches(y), Inches(3.8), Inches(0.4),
             name, size=18, color=WHITE, bold=True)
    # Course & Uni
    add_text(slide, Inches(x), Inches(y + 0.35), Inches(3.8), Inches(0.3),
             f"{course} | {uni}", size=12, color=GRAY)
    # Style
    add_text(slide, Inches(x), Inches(y + 0.65), Inches(3.8), Inches(0.3),
             f"Study: {style}", size=11, color=BRAND_BLUE)

# ===== SLIDE 10: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.4), Inches(1.2), width=Inches(2.5))

add_text(slide, Inches(0.5), Inches(3.8), Inches(12.333), Inches(1),
         "Thank You", size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.6),
         "DEQUAD - Your wellbeing companion for university life", size=22, color=BRAND_BLUE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5),
         "Contact: support@dequad.com", size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/app/DEQUAD_Presentation.pptx"
prs.save(output_path)
file_size = os.path.getsize(output_path)
print(f"Presentation saved to: {output_path}")
print(f"File size: {file_size / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
